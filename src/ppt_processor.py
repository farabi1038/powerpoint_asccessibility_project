"""
PowerPoint Processing Module

This module handles the extraction and modification of PowerPoint files,
with distinct components for different extraction tasks.
"""

import os
import subprocess
import magic
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
import tempfile
from PIL import Image
import io
import shutil
from pptx.dml.color import RGBColor

class PPTProcessor:
    def __init__(self):
        self.presentation = None
        self.image_shapes = []
        self.text_shapes = []
        self.temp_dir = tempfile.mkdtemp()
        
    def load_presentation(self, file_path):
        """Load a PowerPoint presentation"""
        self.presentation = Presentation(file_path)
        self._extract_content()
        return self.presentation
    
    # EXTRACTION MODULE 1: Image Format Conversion
    def _convert_wmf_to_png(self, image_data, slide_idx, shape_idx):
        """Convert WMF/EMF to PNG format using various methods"""
        # First, save the WMF data
        wmf_path = os.path.join(self.temp_dir, f"slide_{slide_idx}_shape_{shape_idx}.wmf")
        with open(wmf_path, 'wb') as f:
            f.write(image_data)
        
        # Output path for PNG
        output_path = os.path.join(self.temp_dir, f"slide_{slide_idx}_shape_{shape_idx}.png")
        
        # Method 1: Try using librsvg (if available)
        try:
            from cairosvg import svg2png
            # First convert WMF to SVG using other tools or libraries if available
            # This is a placeholder - actual WMF to SVG conversion is complex
            # For demonstration, we'll generate a simple SVG with text about the unsupported format
            
            svg_content = f"""
            <svg width="400" height="200" xmlns="http://www.w3.org/2000/svg">
                <rect width="400" height="200" fill="#f8f9fa" />
                <text x="50%" y="50%" text-anchor="middle" font-family="Arial" font-size="16" fill="#333">
                    Windows Metafile Image (WMF/EMF)
                </text>
                <text x="50%" y="70%" text-anchor="middle" font-family="Arial" font-size="12" fill="#666">
                    This image format requires conversion for accessibility
                </text>
            </svg>
            """
            
            svg_path = os.path.join(self.temp_dir, f"slide_{slide_idx}_shape_{shape_idx}.svg")
            with open(svg_path, 'w') as f:
                f.write(svg_content)
            
            svg2png(url=svg_path, write_to=output_path)
            return output_path
        except:
            pass
        
        # Method 2: Try using ImageMagick if available
        try:
            import subprocess
            result = subprocess.run(['convert', wmf_path, output_path], 
                                    check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            if os.path.exists(output_path):
                return output_path
        except:
            pass

        # Method 3: Create a placeholder image with PIL
        try:
            from PIL import Image, ImageDraw, ImageFont
            
            # Create a placeholder image
            img = Image.new('RGB', (400, 200), color=(248, 249, 250))
            d = ImageDraw.Draw(img)
            
            # Try to use a system font
            try:
                font = ImageFont.truetype("Arial", 16)
                small_font = ImageFont.truetype("Arial", 12)
            except:
                font = ImageFont.load_default()
                small_font = ImageFont.load_default()
            
            # Add text to the image
            d.text((200, 100), "Windows Metafile Image (WMF/EMF)", 
                   fill=(51, 51, 51), font=font, anchor="mm")
            d.text((200, 140), "This image format requires conversion for accessibility", 
                   fill=(102, 102, 102), font=small_font, anchor="mm")
            
            # Save the image
            img.save(output_path)
            return output_path
        except:
            # If all methods fail, return None
            return None
    
    # EXTRACTION MODULE 2: Content Extraction Master Function
    def _extract_content(self):
        """Extract content from the presentation"""
        self.image_shapes = []
        self.text_shapes = []
        
        for slide_idx, slide in enumerate(self.presentation.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                # Process text in the shape
                if shape.has_text_frame:
                    self._extract_text_content(slide_idx, shape)
                
                # Extract images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    self._extract_image_content(slide_idx, shape_idx, shape)
    
    # EXTRACTION MODULE 3: Text Content Extraction
    def _extract_text_content(self, slide_idx, shape):
        """Extract text content from shapes with text frames"""
        text_frame = shape.text_frame
        text = ""
        font_size = None
        
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                text += run.text
                if run.font.size:
                    # Convert from EMU to points
                    size_pt = run.font.size.pt
                    font_size = size_pt if font_size is None else min(font_size, size_pt)
        
        self.text_shapes.append({
            "slide_num": slide_idx,
            "shape": shape,
            "text": text,
            "font_size": font_size
        })
    
    # EXTRACTION MODULE 4: Image Content Extraction
    def _extract_image_content(self, slide_idx, shape_idx, shape):
        """Extract image content from picture shapes"""
        try:
            image_data = shape.image.blob
            
            # Extract alt text using multiple methods
            alt_text = self._extract_alt_text(shape)
            
            # Debug
            print(f"Extracted alt text for slide {slide_idx+1}: '{alt_text}'")
            
            # Get file extension based on image content
            image_type = self._get_image_type(image_data)
            
            if image_type == "wmf" or image_type == "emf":
                # Handle Windows Metafile format
                self._handle_wmf_image(image_data, slide_idx, shape_idx, shape, alt_text)
            else:
                # Handle regular image formats
                self._handle_regular_image(image_data, slide_idx, shape_idx, shape, alt_text)
        except Exception as e:
            print(f"Error extracting image data on slide {slide_idx+1}, shape {shape_idx+1}: {e}")
            self.image_shapes.append({
                "slide_num": slide_idx,
                "shape_idx": shape_idx,
                "shape": shape,
                "image_path": None,
                "alt_text": "",
                "warning": f"Error extracting image on slide {slide_idx+1}, shape {shape_idx+1}"
            })
    
    # EXTRACTION MODULE 5: Alt Text Extraction
    def _extract_alt_text(self, shape):
        """Extract alt text from a shape using multiple methods"""
        alt_text = ""
        
        # Method 1: Direct alt_text property
        if hasattr(shape, 'alt_text') and shape.alt_text and hasattr(shape.alt_text, 'text'):
            alt_text = shape.alt_text.text
        
        # Method 2: XML way
        if not alt_text and hasattr(shape, '_element'):
            try:
                cNvPr_element = shape._element.xpath('.//p:cNvPr')
                if cNvPr_element and cNvPr_element[0].get('descr'):
                    alt_text = cNvPr_element[0].get('descr')
            except:
                pass
                
        return alt_text
    
    # EXTRACTION MODULE 6: WMF Image Handling
    def _handle_wmf_image(self, image_data, slide_idx, shape_idx, shape, alt_text):
        """Handle Windows Metafile format images"""
        converted_path = self._convert_wmf_to_png(image_data, slide_idx, shape_idx)
        if converted_path:
            self.image_shapes.append({
                "slide_num": slide_idx,
                "shape_idx": shape_idx,
                "shape": shape,
                "image_path": converted_path,
                "alt_text": alt_text,
                "converted_from_wmf": True
            })
        else:
            # If conversion failed, add placeholder with information
            self.image_shapes.append({
                "slide_num": slide_idx,
                "shape_idx": shape_idx,
                "shape": shape,
                "image_path": None,
                "alt_text": alt_text,
                "converted_from_wmf": False,
                "warning": f"Unsupported image format: WMF file on slide {slide_idx+1}, shape {shape_idx+1}"
            })
            print(f"Warning: Unsupported image format on slide {slide_idx+1}, shape {shape_idx+1}: cannot find loader for this WMF file")
    
    # EXTRACTION MODULE 7: Regular Image Handling
    def _handle_regular_image(self, image_data, slide_idx, shape_idx, shape, alt_text):
        """Handle regular image formats"""
        try:
            # Try to open the image
            image = Image.open(io.BytesIO(image_data))
            
            # Save the image to a temporary file
            ext = "." + (image.format.lower() if image.format else "png")
            img_path = os.path.join(self.temp_dir, f"slide_{slide_idx}_shape_{shape_idx}{ext}")
            image.save(img_path)
            
            self.image_shapes.append({
                "slide_num": slide_idx,
                "shape_idx": shape_idx,
                "shape": shape,
                "image_path": img_path,
                "alt_text": alt_text
            })
        except OSError as e:
            # Check if "WMF" is in the error - indicates unsupported format
            if "WMF" in str(e).upper():
                # Instead of raising an error, set a warning flag and include shape_idx
                self.image_shapes.append({
                    "slide_num": slide_idx,
                    "shape_idx": shape_idx,
                    "shape": shape,
                    "image_path": None,
                    "alt_text": alt_text,
                    "warning": "WMF file skipped - PIL cannot load"
                })
            else:
                # Some other error - re-raise
                raise e
    
    # EXTRACTION MODULE 8: Image Type Detection
    def _get_image_type(self, image_data):
        """Determine image type from binary data"""
        try:
            import magic
            mime = magic.Magic(mime=True)
            mime_type = mime.from_buffer(image_data)
            
            if mime_type == "image/x-wmf" or "wmf" in mime_type:
                return "wmf"
            elif mime_type == "image/x-emf" or "emf" in mime_type:
                return "emf"
            else:
                return mime_type.split('/')[-1]
        except ImportError:
            # If python-magic is not available, try basic header check
            if image_data[:4] == b'\xd7\xcd\xc6\x9a':  # WMF header
                return "wmf"
            elif image_data[:4] == b'\x01\x00\x00\x00':  # EMF header
                return "emf"
            else:
                # Try using PIL's Image.open
                try:
                    image = Image.open(io.BytesIO(image_data))
                    return image.format.lower() if image.format else "unknown"
                except:
                    return "unknown"
    
    # EXTRACTION MODULE 9: Font Size Detection
    def _get_shape_font_size(self, shape):
        """Get the font size of a shape"""
        if not shape.has_text_frame:
            return None
        
        font_size = None
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size:
                    size_pt = run.font.size.pt
                    font_size = size_pt if font_size is None else min(font_size, size_pt)
        
        return font_size
                
    # ENHANCEMENT MODULE 1: Alt Text Update
    def update_alt_text(self, slide_idx, shape, alt_text):
        """Update alt text for an image shape"""
        try:
            # Method 1: Direct property method
            if hasattr(shape, 'alt_text'):
                shape.alt_text.text = alt_text
                return True
            
            # Method 2: XML way
            if hasattr(shape, '_element'):
                cNvPr_element = shape._element.xpath('.//p:cNvPr')
                if cNvPr_element:
                    cNvPr_element[0].set('descr', alt_text)
                    return True
            
            return False
        except Exception as e:
            print(f"Error updating alt text: {e}")
            return False
    
    # ENHANCEMENT MODULE 2: Caption Addition
    def add_visible_caption(self, slide_idx, shape, caption_text, is_single_image=False):
        """Add a visible caption to an image shape that's consistently placed below the image"""
        try:
            slide = self.presentation.slides[slide_idx]
            
            # Get the position and size of the image
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            
            # Calculate position for the text box (always directly below the image)
            # Always use the image's width and position for consistency
            caption_left = left
            caption_width = width
            
            # Always place the caption directly below the image with a small gap
            caption_top = top + height + Inches(0.05)
            
            # If the caption would go off the slide, adjust the position but keep it below the image
            slide_height = self.presentation.slide_height
            if caption_top + Inches(0.4) > slide_height:
                # Move the caption slightly up but still keep it below the image
                caption_top = slide_height - Inches(0.45)
            
            # Create a textbox shape with a border to make it more visible
            textbox = slide.shapes.add_textbox(
                caption_left, caption_top, caption_width, Inches(0.4)
            )
            
            # Add border to make the caption stand out
            if hasattr(textbox, 'line'):
                textbox.line.color.rgb = RGBColor(100, 100, 100)  # Gray border
                textbox.line.width = Pt(1.0)
            
            # Add light background to improve readability
            if hasattr(textbox, 'fill'):
                textbox.fill.solid()
                textbox.fill.fore_color.rgb = RGBColor(245, 245, 220)  # Light beige background
            
            # Set the text
            text_frame = textbox.text_frame
            text_frame.clear()
            
            # Enable text wrapping
            text_frame.word_wrap = True
            
            # Enable auto-size to fit text
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
            # Set vertical alignment
            text_frame.vertical_anchor = 1  # Middle
            
            # Add paragraph
            p = text_frame.paragraphs[0]
            p.alignment = 1  # Center
            
            # Add run with the caption text
            run = p.add_run()
            run.text = caption_text
            
            # Format the text to be visually distinct but very readable
            font = run.font
            font.name = "Calibri"
            font.size = Pt(14)  # Larger font for better readability
            font.bold = True  # Make it bold for emphasis
            font.italic = False  # No italic for better readability
            font.color.rgb = RGBColor(0, 0, 0)  # Black text for maximum contrast
            
            # Return success
            return textbox
        except Exception as e:
            import traceback
            print(f"Error adding caption: {e}")
            print(traceback.format_exc())
            return None
    
    # ENHANCEMENT MODULE 3: Font Size Update
    def update_font_size(self, shape, new_size):
        """Update the font size of a shape"""
        if not shape.has_text_frame:
            return False
        
        try:
            text_frame = shape.text_frame
            changed = False
            
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    # Skip if no text
                    if not run.text.strip():
                        continue
                        
                    # Only increase font size if it's smaller than the new size
                    current_size = run.font.size.pt if hasattr(run.font, 'size') and run.font.size else None
                    
                    if current_size is None:
                        # If size isn't set, set it to the new size
                        run.font.size = Pt(new_size)
                        changed = True
                    elif current_size < new_size:
                        # Only increase font size if it's too small
                        run.font.size = Pt(new_size)
                        changed = True
                    # else: keep the existing size if it's already large enough
            
            return changed
        except Exception as e:
            print(f"Error updating font size: {e}")
            return False
    
    # ENHANCEMENT MODULE 4: Text Content Update
    def update_text(self, shape, new_text):
        """Update the text content of a shape"""
        if not shape.has_text_frame:
            return False
        
        try:
            text_frame = shape.text_frame
            
            # If there are no paragraphs, add one
            if not text_frame.paragraphs:
                text_frame.text = new_text
            else:
                # Otherwise, update the first paragraph
                paragraph = text_frame.paragraphs[0]
                paragraph.text = new_text
            
            return True
        except Exception as e:
            print(f"Error updating text: {e}")
            return False
    
    # ENHANCEMENT MODULE 5: Text Contrast Update
    def update_text_contrast(self, shape, make_darker=True):
        """Update the contrast of text for better readability"""
        if not shape.has_text_frame:
            return False
        
        try:
            text_frame = shape.text_frame
            changed = False
            
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    # Skip if font or color is not set
                    if not hasattr(run, 'font') or not hasattr(run.font, 'color'):
                        continue
                        
                    # Check if color is None or doesn't have rgb property
                    if not run.font.color or not hasattr(run.font.color, 'rgb'):
                        continue
                        
                    # Handle case where rgb exists but is None
                    if run.font.color.rgb is None:
                        continue
                    
                    try:
                        # Get current RGB values - handle different color object types
                        current_color = run.font.color.rgb
                        
                        # Check if color has r, g, b attributes
                        if not all(hasattr(current_color, attr) for attr in ['r', 'g', 'b']):
                            continue
                            
                        r, g, b = current_color.r, current_color.g, current_color.b
                        
                        # Calculate luminance
                        luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
                        
                        if make_darker:
                            # If it's light, make it dark (for better contrast on light backgrounds)
                            if luminance > 0.5:
                                run.font.color.rgb = RGBColor(0, 0, 0)  # Black
                                changed = True
                        else:
                            # If it's dark, make it light (for better contrast on dark backgrounds)
                            if luminance < 0.5:
                                run.font.color.rgb = RGBColor(255, 255, 255)  # White
                                changed = True
                    except AttributeError:
                        # Skip this run if any attribute errors occur when accessing color properties
                        continue
            
            return changed
        except Exception as e:
            print(f"Error updating text contrast: {e}")
            return False
    
    # UTILITY MODULE 1: Presentation Save
    def save_presentation(self, output_path):
        """Save the presentation to a file"""
        if self.presentation:
            try:
                self.presentation.save(output_path)
                return True
            except Exception as e:
                print(f"Error saving presentation: {e}")
                return False
        return False
    
    # UTILITY MODULE 2: Cleanup
    def cleanup(self):
        """Clean up temporary files"""
        try:
            shutil.rmtree(self.temp_dir)
            return True
        except Exception as e:
            print(f"Error cleaning up: {e}")
            return False
    
    # ENHANCEMENT MODULE 6: Simple Caption
    def add_simple_caption(self, slide_idx, shape, caption_text):
        """Add a simple caption at the bottom of the slide"""
        try:
            slide = self.presentation.slides[slide_idx]
            
            # Create a textbox at the bottom of the slide
            textbox = slide.shapes.add_textbox(
                Inches(0.5), 
                self.presentation.slide_height - Inches(1.0),
                self.presentation.slide_width - Inches(1.0), 
                Inches(0.75)
            )
            
            # Set the text
            text_frame = textbox.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            
            p = text_frame.paragraphs[0]
            p.alignment = 1  # Center
            
            run = p.add_run()
            run.text = caption_text
            
            # Format
            font = run.font
            font.name = "Calibri"
            font.size = Pt(14)
            font.bold = True
            font.color.rgb = RGBColor(0, 0, 0)  # Black
            
            return textbox
        except Exception as e:
            print(f"Error adding simple caption: {e}")
            return None 
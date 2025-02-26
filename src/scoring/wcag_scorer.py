# src/scoring/scorer.py
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from src.llm.llm_processor import LLMProcessor
from pptx.dml.color import RGBColor

class WCAGScorer:
    """
    Analyzes a PowerPoint presentation for accessibility issues based on WCAG 2.1 guidelines.
    Applies fixes using an LLM (if configured) and computes an overall score.

    Checks include:
      1. Slide Title Presence
      2. Image Alt Text
      3. Table Header Rows
      4. Empty Text Boxes
      5. Color Contrast (rudimentary check - placeholder for more advanced logic)
      6. Font Size (minimum recommended size check)
      7. Link Text Clarity
    """

    def __init__(self, llm_processor: LLMProcessor):
        """
        :param llm_processor: An instance of LLMProcessor for generating or refining alt text, 
                              slide titles, or other textual improvements.
        """
        self.llm_processor = llm_processor
        self.initial_score = 100
        self.score = self.initial_score
        self.issues_log = []

    def analyze_presentation(self, pptx_stream):
        """
        Processes each slide in the presentation.

        :param pptx_stream: A file-like object containing the .pptx data.
        :return: (Presentation object, final score, list of issues/fixes)
        """
        prs = Presentation(pptx_stream)
        for slide_index, slide in enumerate(prs.slides, start=1):
            self.analyze_slide(slide, slide_index)

        # Ensure score doesn't go below 0
        self.score = max(0, self.score)

        return prs, self.score, self.issues_log

    def analyze_slide(self, slide, slide_index: int):
        """
        Analyzes a single slide for various accessibility issues.

        :param slide: The slide object to analyze.
        :param slide_index: The current slide number (1-based).
        """
        slide_title = self.check_slide_title(slide, slide_index)
        self.check_images(slide, slide_index)
        self.check_tables(slide, slide_index)
        self.check_empty_text_boxes(slide, slide_index)
        self.check_color_contrast(slide, slide_index)
        self.check_font_size(slide, slide_index)
        self.check_link_texts(slide, slide_index)

    def check_slide_title(self, slide, slide_index):
        """
        Ensures each slide has a title. Deducts points and logs issue if missing.
        Uses shape naming or text heuristics to detect title shapes.
        """
        slide_title = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                # Attempt to identify a 'title' shape from its name or placeholder usage
                if "title" in shape.name.lower() and text:
                    slide_title = text
                    break

        if not slide_title:
            self.score -= 5
            self.issues_log.append(f"Slide {slide_index}: Missing slide title (-5 points).")

        return slide_title

    def check_images(self, slide, slide_index):
        """
        Checks for images (pictures) without alt text. Deducts points if missing.
        Optionally calls the LLM to generate alt text (if desired) as a fix.
        """
        for shape in slide.shapes:
            # Only check pictures
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                alt_text = self._get_alt_text(shape)
                if not alt_text:
                    self.score -= 10
                    self.issues_log.append(
                        f"Slide {slide_index}: Image missing alt text (-10 points)."
                    )
                    # Optionally generate and apply alt text via LLM
                    # generated_alt_text = self.llm_processor.generate_alt_text_from_image(...)
                    # self._set_alt_text(shape, generated_alt_text)

    def check_tables(self, slide, slide_index):
        """
        Checks for tables without a designated header row. Deducts points if missing.
        """
        for shape in slide.shapes:
            if hasattr(shape, "has_table") and shape.has_table:
                table = shape.table
                if table.rows:
                    # We assume the first row is the header row. Check if it has meaningful text.
                    header_found = any(cell.text.strip() for cell in table.rows[0].cells)
                    if not header_found:
                        self.score -= 5
                        self.issues_log.append(
                            f"Slide {slide_index}: Table missing header row description (-5 points)."
                        )

    def check_empty_text_boxes(self, slide, slide_index):
        """
        Flags empty text boxes to detect potential missing or incomplete descriptions.
        """
        for shape in slide.shapes:
            if shape.has_text_frame:
                if not shape.text.strip():
                    self.score -= 2
                    self.issues_log.append(
                        f"Slide {slide_index}: Text shape is empty, indicating missing description (-2 points)."
                    )

    def check_color_contrast(self, slide, slide_index):
        """
        Checks for text that may have poor contrast. This is a very basic placeholder:
        - Assumes a white slide background.
        - Retrieves the color from the first run in each text frame (if it exists).
        - Deducts points if the color is "too light" (placeholder logic).
        """

        def is_contrast_poor(rgb_color):
            # Example rule: if all channels are above 200, it's considered "too light"
            return (rgb_color[0] > 200 and 
                    rgb_color[1] > 200 and 
                    rgb_color[2] > 200)

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # Safely check if there's a color object and if that color object has an .rgb attribute
                        if run.font.color is not None:
                            # Some runs may have an _NoneColor object, so we need to ensure it's not that type
                            color_obj = run.font.color
                            # If color_obj is an RGBColor instance, we can safely check .rgb
                            if hasattr(color_obj, "rgb") and color_obj.rgb:
                                color_rgb = color_obj.rgb
                                # color_rgb is a pptx.dml.color.RGBColor object; we can treat it like a tuple
                                if is_contrast_poor(color_rgb):
                                    self.score -= 3
                                    self.issues_log.append(
                                        f"Slide {slide_index}: Text color may lack contrast (-3 points)."
                                    )
                                    # Consider whether you want to break after first detection for the shape
                                    # break

    def check_font_size(self, slide, slide_index):
        """
        Checks the font size in text frames. Deducts points if smaller than ~18 pt 
        (common guideline for presentations to ensure readability).
        """
        MIN_FONT_SIZE = 18.0

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size:
                            # font.size is in EMU; to convert to Pt: 1 pt = 12700 EMU
                            pt_size = run.font.size.pt
                            if pt_size and pt_size < MIN_FONT_SIZE:
                                self.score -= 2
                                self.issues_log.append(
                                    f"Slide {slide_index}: Text font size ({pt_size:.1f} pt) is below recommended minimum ({MIN_FONT_SIZE} pt) (-2 points)."
                                )

    def check_link_texts(self, slide, slide_index):
        """
        Checks for link texts that may be unclear (e.g., "click here" or raw URLs).
        Deducts points if found, since descriptive link text is a WCAG recommendation.
        """
        unclear_phrases = ["click here", "more info", "read more"]

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.lower()
                    # Check if text is purely a URL or has unclear phrases
                    # This is a simple heuristic; a real solution would parse runs for actual hyperlink data.
                    if "http://" in text or "https://" in text:
                        # Potential raw URL
                        self.score -= 2
                        self.issues_log.append(
                            f"Slide {slide_index}: Link text might be a raw URL, consider descriptive text (-2 points)."
                        )
                    for phrase in unclear_phrases:
                        if phrase in text:
                            self.score -= 2
                            self.issues_log.append(
                                f"Slide {slide_index}: Link text '{phrase}' is unclear (-2 points)."
                            )

    # -------------------------------------------------------------
    # Helper Methods
    # -------------------------------------------------------------
    def _get_alt_text(self, shape):
        """
        Retrieves alt text from the shape if it exists.
        """
        # python-pptx does not expose alt text directly via a property.
        # Usually, shape._element.xpath can read it from p:cNvPr@descr or p:cNvPr@title.
        if not hasattr(shape, '_element'):
            return None

        descr_list = shape._element.xpath('.//p:cNvPr/@descr')
        if descr_list:
            return descr_list[0]  # Return first matching alt text
        return None

    def _set_alt_text(self, shape, text):
        """
        Sets or updates the alt text for the given shape.
        Example usage after generating alt text with LLM.
        """
        if hasattr(shape, '_element'):
            cNvPr = shape._element.xpath('.//p:cNvPr')
            if cNvPr:
                cNvPr[0].set('descr', text)

import streamlit as st
from pptx import Presentation
from openai import OpenAI

st.set_page_config(page_title="PowerPoint LLM Processor", layout="wide")

st.title("📊 PowerPoint LLM Processor with Drag & Drop")

st.write("Upload your PowerPoint file and let the LLM analyze it dynamically!")

uploaded_file = st.file_uploader("**Drag and Drop or Click to Upload**", type=["pptx"], accept_multiple_files=False)

@st.cache_data
def extract_text_from_pptx(file):
    presentation = Presentation(file)
    text = []
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        text.append(" ".join(slide_text))
    return text

if uploaded_file is not None:
    st.success("✅ File uploaded successfully!")
    ppt_text = extract_text_from_pptx(uploaded_file)

    # Display extracted text
    st.subheader("Extracted Content from Slides:")
    for i, slide_text in enumerate(ppt_text, start=1):
        st.markdown(f"### Slide {i}")
        st.write(slide_text)

    # Placeholder for LLM API integration
    if st.button("Process with LLM"):
        with st.spinner("🤖 Processing with LLM..."):
            # Call your LLM API here (e.g., OpenAI API)
            llm_response = f"Here would be the LLM response based on {len(ppt_text)} slides of content."
            st.success("🎉 Processing complete!")
            st.text_area("LLM Output", llm_response, height=300)
else:
    st.info("Upload a PowerPoint file to get started!")